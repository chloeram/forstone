# slides_parser.py

import pdfplumber
from pptx import Presentation
import re
from src.content_models import ContentItem

def clean_text(text: str) -> str:
    text = re.sub(r"^\s*[\•\-\*]\s*", "", text, flags=re.MULTILINE)  # Remove bullets at the start of lines
    text = re.sub(r"\s+", " ", text).strip()  # Normalize spaces
    return text

# Parses slides in PDF format
def parse_pdf_slides(filepath: str) -> list[dict]: 

    slides = []

    # Opening PDF file
    with pdfplumber.open(filepath) as pdf:
        # Iterate through pages (slides) of the PDF document to extract content, starting at index 1 for "slide 1"
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or "" # Extracts text from page if any, else None
            slide_text = clean_text(text.strip())

            metadata = {"slide_number": str(i)}

            item = ContentItem(
                source="slide",
                text=slide_text,
                metadata=metadata
            )

            slides.append(item)
            
    return slides

# Parses slides in PPT/PPTX format
def parse_pptx_slides(filepath: str) -> list[dict]:
    
    slides = []
    ppt = Presentation(filepath)

    for i, slide in enumerate(ppt.slides, start = 1): # Iterate through slides, count slides for ref

        title = ""
        slide_text = []

        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.has_text_frame:
                continue

            text = clean_text(shape.text.strip())  # Clean bullet points
            if shape.placeholder_format and shape.placeholder_format.type == 1 and not title:
                title = text  # Extract slide title
            else:
                slide_text.append(text)

        combined_text = title + "\n" + "\n".join(slide_text)
        metadata = {"slide_number": str(i)}
        
        if title:
            metadata["title"] = title
        
        item = ContentItem(
            source="slide",
            text=combined_text,
            metadata=metadata
        )

        slides.append(item)

    return slides
